"""
Process different file types for document ingestion.
"""
import os
import re
import unicodedata
from typing import Dict, Any
from pathlib import Path

# Optional imports
try:
    import docx
except ImportError:
    docx = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import textract
except ImportError:
    textract = None

try:
    import extract_msg
except ImportError:
    extract_msg = None

try:
    import markdown
    from bs4 import BeautifulSoup
    markdown_available = True
except ImportError:
    markdown_available = False

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None

from bs4 import BeautifulSoup

# Replace with your own logger if needed
import logging
logger = logging.getLogger(__name__)

class FileProcessor:
    """Process different file types for text extraction."""

    SUPPORTED_EXTENSIONS = {
        '.pdf', '.docx', '.doc', '.txt', '.md', '.csv', '.xls', '.xlsx',
        '.pptx', '.html', '.htm', '.msg', '.jpg', '.jpeg', '.png'
    }
    MAX_FILE_SIZE_MB = 20

    @staticmethod
    def process_file(file_path: Path) -> Dict[str, Any]:
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return FileProcessor._error_doc(file_path, "FILE NOT FOUND")

        if file_path.stat().st_size > FileProcessor.MAX_FILE_SIZE_MB * 1024 * 1024:
            logger.warning(f"File too large: {file_path.name}")
            return FileProcessor._error_doc(file_path, "FILE TOO LARGE")

        file_extension = file_path.suffix.lower()
        file_name = file_path.name
        title = FileProcessor._generate_title(file_name, file_extension)

        handler = FileProcessor._get_handler(file_extension)

        try:
            content = handler(file_path)
            content = unicodedata.normalize('NFKC', content or "").strip()
            if not content:
                content = f"[EMPTY CONTENT: {file_path.name}]"
            elif len(content) < 50:
                logger.warning(f"Very short content extracted from {file_path} ({len(content)} characters)")

            logger.info(f"Processed {file_extension.upper()} file: {file_name} ({len(content)} characters)")
            return {
                "title": title,
                "content": content,
                "source_file": str(file_path),
                "file_type": file_extension.lstrip('.')
            }

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return FileProcessor._error_doc(file_path, str(e), title)

    @staticmethod
    def _get_handler(extension: str):
        return {
            '.pdf': FileProcessor.extract_from_pdf,
            '.docx': FileProcessor.extract_from_docx,
            '.doc': FileProcessor.extract_from_doc,
            '.txt': FileProcessor.extract_from_txt,
            '.md': FileProcessor.extract_from_markdown,
            '.csv': FileProcessor.extract_from_csv,
            '.xls': FileProcessor.extract_from_excel,
            '.xlsx': FileProcessor.extract_from_excel,
            '.pptx': FileProcessor.extract_from_pptx,
            '.html': FileProcessor.extract_from_html,
            '.htm': FileProcessor.extract_from_html,
            '.msg': FileProcessor.extract_from_msg,
            '.jpg': FileProcessor.extract_from_image,
            '.jpeg': FileProcessor.extract_from_image,
            '.png': FileProcessor.extract_from_image
        }.get(extension, FileProcessor.extract_from_txt)

    @staticmethod
    def _generate_title(file_name: str, file_extension: str) -> str:
        title = file_name.replace(file_extension, '').replace('-', ' ').replace('_', ' ')
        return re.sub(r'\s+', ' ', title).strip().title()

    @staticmethod
    def _error_doc(file_path: Path, error: str, title: str = None) -> Dict[str, Any]:
        return {
            "title": title or f"Error: {file_path.name}",
            "content": f"[{error}: {file_path.name}]",
            "source_file": str(file_path),
            "file_type": "error"
        }

    @staticmethod
    def extract_from_pdf(file_path: Path) -> str:
        if fitz is None:
            return "[PDF support not installed]"
        with fitz.open(file_path) as pdf:
            return " ".join(page.get_text() for page in pdf)

    @staticmethod
    def extract_from_docx(file_path: Path) -> str:
        if docx is None:
            return "[DOCX support not installed]"
        return "\n".join(p.text for p in docx.Document(file_path).paragraphs)

    @staticmethod
    def extract_from_doc(file_path: Path) -> str:
        if textract is None:
            return "[DOC support not installed]"
        return textract.process(str(file_path)).decode()

    @staticmethod
    def extract_from_txt(file_path: Path) -> str:
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        return open(file_path, 'rb').read().decode('utf-8', errors='replace')

    @staticmethod
    def extract_from_markdown(file_path: Path) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            md_text = f.read()
        if not markdown_available:
            return md_text
        html = markdown.markdown(md_text)
        return BeautifulSoup(html, 'html.parser').get_text()

    @staticmethod
    def extract_from_csv(file_path: Path) -> str:
        if pd is None:
            return "[Pandas not installed for CSV processing]"
        return pd.read_csv(file_path).to_string(index=False)

    @staticmethod
    def extract_from_excel(file_path: Path) -> str:
        if pd is None:
            return "[Pandas not installed for Excel processing]"
        return pd.read_excel(file_path).to_string(index=False)

    @staticmethod
    def extract_from_pptx(file_path: Path) -> str:
        if Presentation is None:
            return "[PPTX support not installed]"
        prs = Presentation(file_path)
        return " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

    @staticmethod
    def extract_from_html(file_path: Path) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        return soup.get_text()

    @staticmethod
    def extract_from_msg(file_path: Path) -> str:
        if extract_msg is None:
            return "[MSG support not installed]"
        msg = extract_msg.Message(str(file_path))
        return f"{msg.subject}\n{msg.body}"

    @staticmethod
    def extract_from_image(file_path: Path) -> str:
        if Image is None or pytesseract is None:
            return "[OCR support not installed]"
        return pytesseract.image_to_string(Image.open(file_path))